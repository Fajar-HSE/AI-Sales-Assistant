#!/usr/bin/env python3
"""
Sales AI Assistant — FastAPI Backend v0.5 (Multi-User)
Gateway WA : Fonnte (api.fonnte.com)
LLM       : Groq (llama-3.3-70b-versatile, JSON mode) + fallback rule-based
DB        : SQLite (multi-user: users, customers, chats, knowledge_base, products)
Auth      : JWT per-user + RBAC (admin = super-admin, user = own data only)
Endpoints : /health, /ws, /webhook/fonnte (+ /webhook/fonte alias),
            /api/v1/auth/login, /api/v1/me, /api/v1/admin/users,
            /api/v1/messages/send, /api/v1/assessment/analyze,
            /api/v1/reply/generate, /api/v1/customers,
            /api/v1/products/*, /api/v1/upload, /api/v1/knowledge/*
Security  : webhook token, JWT/API-token auth, CORS restricted, rate limit
"""
import os, json, hmac, hashlib, uuid, datetime, logging, io, time, sqlite3, threading
import httpx

from dotenv import load_dotenv
from fastapi import FastAPI, Request, HTTPException, Depends, WebSocket, WebSocketDisconnect, UploadFile, File, Form, Header
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse

try:
    import jwt as pyjwt
except Exception:
    pyjwt = None

load_dotenv()  # baca .env
logging.basicConfig(level=logging.INFO, format="%(asctime)s %(message)s")
log = logging.getLogger("sai")

# ---------- CONFIG ----------
FONNTE_TOKEN      = os.getenv("FONNTE_TOKEN", "")
FONNTE_FROM       = os.getenv("FONNTE_FROM_NUMBER", "6289876543210")
GROQ_API_KEY      = os.getenv("GROQ_API_KEY", "")
GROQ_MODEL        = os.getenv("GROQ_MODEL", "llama-3.3-70b-versatile")

# Webhook — jika diisi, /webhook/* hanya menerima token yang cocok
WEBHOOK_SECRET = os.getenv("WEBHOOK_SECRET", "")

# Auth — JWT (login) dan/atau static API token. Jika keduanya kosong => dev mode terbuka.
JWT_SECRET_KEY  = os.getenv("JWT_SECRET_KEY", "")
JWT_EXPIRATION  = int(os.getenv("JWT_EXPIRATION", "86400") or "86400")
API_TOKEN       = os.getenv("API_TOKEN", "")
ADMIN_USER      = os.getenv("ADMIN_USER", "admin")
ADMIN_PASSWORD  = os.getenv("ADMIN_PASSWORD", "")
AUTH_ENABLED    = bool(JWT_SECRET_KEY or API_TOKEN)

# Rate limit (in-memory, per IP/key)
RATE_WINDOW = int(os.getenv("RATE_WINDOW", "60"))
RATE_MAX    = int(os.getenv("RATE_MAX", "60"))

# Frontend (index.html) dilayani di "/" — default relatif ke file backend ini
FRONTEND_FILE = os.getenv("FRONTEND_FILE", os.path.join(os.path.dirname(os.path.abspath(__file__)), "index.html"))

# CORS — frontend dilayani same-origin via /salesai/, sehingga origin ketat aman.
cors_origins = [o.strip() for o in os.getenv("CORS_ORIGINS", "").split(",") if o.strip()] or \
               ["http://localhost:8000", "http://127.0.0.1:8000"]

app = FastAPI(title="Sales AI Assistant API", version="0.5.0")
app.add_middleware(CORSMiddleware,
                   allow_origins=cors_origins,
                   allow_methods=["*"],
                   allow_headers=["*"])

# Supabase (opsional, kompatibilitas env lama; store utama sekarang SQLite)
SUPABASE_URL = os.getenv("SUPABASE_URL", "")
SUPABASE_KEY = os.getenv("SUPABASE_KEY", "")
SUPABASE_SERVICE_KEY = os.getenv("SUPABASE_SERVICE_KEY", "")
HAS_SUPABASE = bool(SUPABASE_URL and (SUPABASE_KEY or SUPABASE_SERVICE_KEY))
try:
    from supabase import create_client
    supabase_admin = create_client(SUPABASE_URL, SUPABASE_SERVICE_KEY or SUPABASE_KEY)
    log.warning("[Supabase] Connected to %s", SUPABASE_URL)
except Exception as e:
    supabase_admin = None
    log.warning("[Supabase] Not available (SQLite digunakan sebagai store utama): %s", e)

# ---------- SQLITE STORE ----------
DB_PATH = os.getenv("DB_PATH", os.path.join(os.path.dirname(os.path.abspath(__file__)), "data", "app.db"))
os.makedirs(os.path.dirname(DB_PATH), exist_ok=True)
_DB_LOCK = threading.Lock()

def conn():
    c = sqlite3.connect(DB_PATH, check_same_thread=False)
    c.row_factory = sqlite3.Row
    return c

def db_exec(sql, params=()):
    with _DB_LOCK:
        c = conn()
        try:
            c.execute(sql, params)
            c.commit()
        finally:
            c.close()

def db_query(sql, params=()):
    c = conn()
    try:
        rows = c.execute(sql, params).fetchall()
    finally:
        c.close()
    return [dict(r) for r in rows]

def init_db():
    db_exec("""CREATE TABLE IF NOT EXISTS users (
        id TEXT PRIMARY KEY,
        username TEXT UNIQUE NOT NULL,
        password_hash TEXT NOT NULL,
        role TEXT NOT NULL DEFAULT 'user',
        display_name TEXT DEFAULT '',
        fonnte_token TEXT DEFAULT '',
        fonnte_from_number TEXT DEFAULT '',
        groq_api_key TEXT DEFAULT '',
        is_active INTEGER DEFAULT 1,
        created_at TEXT,
        updated_at TEXT
    )""")
    db_exec("""CREATE TABLE IF NOT EXISTS customers (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        owner_id TEXT NOT NULL,
        phone TEXT NOT NULL,
        name TEXT DEFAULT '',
        last_score INTEGER DEFAULT 0,
        category TEXT DEFAULT 'Cold Lead',
        badge TEXT DEFAULT '🔴',
        product TEXT DEFAULT 'Umum',
        last_message TEXT DEFAULT '',
        unread INTEGER DEFAULT 0,
        created_at TEXT,
        last_updated TEXT,
        UNIQUE(owner_id, phone)
    )""")
    db_exec("""CREATE TABLE IF NOT EXISTS chats (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        owner_id TEXT NOT NULL,
        phone TEXT NOT NULL,
        text TEXT DEFAULT '',
        direction TEXT DEFAULT 'in',
        timestamp TEXT,
        lead_score INTEGER,
        category TEXT,
        intent_label TEXT,
        sender_name TEXT
    )""")
    db_exec("""CREATE TABLE IF NOT EXISTS knowledge_base (
        id TEXT PRIMARY KEY,
        owner_id TEXT NOT NULL,
        category TEXT DEFAULT 'Umum',
        name TEXT DEFAULT '',
        filename TEXT DEFAULT '',
        file_url TEXT DEFAULT '',
        kb_text TEXT DEFAULT '',
        chunks TEXT DEFAULT '[]',
        chunk_count INTEGER DEFAULT 0,
        uploaded_at TEXT
    )""")
    db_exec("""CREATE TABLE IF NOT EXISTS products (
        id TEXT PRIMARY KEY,
        owner_id TEXT NOT NULL,
        name TEXT DEFAULT '',
        category TEXT DEFAULT 'Umum',
        description TEXT DEFAULT '',
        price_range TEXT DEFAULT '',
        duration TEXT DEFAULT '',
        kb_text TEXT DEFAULT '',
        created_at TEXT
    )""")
    db_exec("CREATE INDEX IF NOT EXISTS idx_cust_owner ON customers(owner_id)")
    db_exec("CREATE INDEX IF NOT EXISTS idx_chats_owner_phone ON chats(owner_id, phone)")
    db_exec("CREATE INDEX IF NOT EXISTS idx_kb_owner ON knowledge_base(owner_id)")

# ---------- PASSWORD HASHING ----------
def hash_password(pw: str) -> str:
    salt = os.urandom(16).hex()
    dk = hashlib.pbkdf2_hmac("sha256", pw.encode(), bytes.fromhex(salt), 100000).hex()
    return f"pbkdf2_sha256$100000${salt}${dk}"

def verify_password(pw: str, stored: str) -> bool:
    try:
        _, it, salt, dk = stored.split("$")
        calc = hashlib.pbkdf2_hmac("sha256", pw.encode(), bytes.fromhex(salt), int(it)).hex()
        return hmac.compare_digest(calc, dk)
    except Exception:
        return False

# ---------- USER HELPERS ----------
def _row_to_user(r: dict) -> dict:
    return {
        "id": r["id"], "username": r["username"], "role": r["role"],
        "display_name": r.get("display_name") or r["username"],
        "fonnte_token": r.get("fonnte_token") or "",
        "fonnte_from_number": r.get("fonnte_from_number") or "",
        "groq_api_key": r.get("groq_api_key") or "",
        "is_active": bool(r.get("is_active", 1)),
        "created_at": r.get("created_at"),
    }

def get_user_by_id(uid: str):
    rows = db_query("SELECT * FROM users WHERE id=?", (uid,))
    return _row_to_user(rows[0]) if rows else None

def get_user_by_username(uname: str):
    rows = db_query("SELECT * FROM users WHERE username=?", (uname,))
    return _row_to_user(rows[0]) if rows else None

def list_users():
    rows = db_query("SELECT * FROM users ORDER BY created_at ASC")
    return [_row_to_user(r) for r in rows]

def create_user(uname, pw, role="user", display_name="", fonnte_token="", fonnte_from="", groq_key=""):
    uid = str(uuid.uuid4())
    now = datetime.datetime.now(datetime.timezone.utc).isoformat()
    db_exec(
        "INSERT INTO users (id,username,password_hash,role,display_name,fonnte_token,fonnte_from_number,groq_api_key,is_active,created_at,updated_at) "
        "VALUES (?,?,?,?,?,?,?,?,1,?,?)",
        (uid, uname, hash_password(pw), role, display_name or uname, fonnte_token, fonnte_from, groq_key, now, now))
    return get_user_by_id(uid)

def update_user(uid, **fields):
    allowed = {"role","display_name","fonnte_token","fonnte_from_number","groq_api_key","is_active"}
    sets, params = [], []
    for k, v in fields.items():
        if k in allowed:
            sets.append(f"{k}=?")
            params.append(v)
    if not sets:
        return get_user_by_id(uid)
    sets.append("updated_at=?")
    params.append(datetime.datetime.now(datetime.timezone.utc).isoformat())
    params.append(uid)
    db_exec(f"UPDATE users SET {','.join(sets)} WHERE id=?", tuple(params))
    return get_user_by_id(uid)

def delete_user(uid):
    with _DB_LOCK:
        c = conn()
        try:
            c.execute("DELETE FROM users WHERE id=?", (uid,))
            c.execute("DELETE FROM customers WHERE owner_id=?", (uid,))
            c.execute("DELETE FROM chats WHERE owner_id=?", (uid,))
            c.execute("DELETE FROM knowledge_base WHERE owner_id=?", (uid,))
            c.execute("DELETE FROM products WHERE owner_id=?", (uid,))
            c.commit()
        finally:
            c.close()

def set_password(uid, pw):
    db_exec("UPDATE users SET password_hash=?, updated_at=? WHERE id=?",
            (hash_password(pw), datetime.datetime.now(datetime.timezone.utc).isoformat(), uid))

def seed_admin():
    rows = db_query("SELECT id FROM users LIMIT 1")
    if rows:
        return
    uname = ADMIN_USER or "admin"
    pw = ADMIN_PASSWORD or "admin123"
    u = create_user(uname, pw, "admin", display_name=uname)
    log.warning("[Seed] Admin dibuat: %s / (password dari env ADMIN_PASSWORD)", uname)

init_db()
seed_admin()

# ---------- MOCK DATA (fallback in-memory untuk KB saat no-DB — tidak dipakai, SQLite selalu ada) ----------
KB_STORE = []
PRODUCT_CATEGORIES = ["BNSP", "Kemnaker RI", "Reguler", "Umum"]
KB_CHUNK_SIZE = 500
KB_CHUNK_OVERLAP = 80

# ---------- HELPERS ----------
def _to_int(v, default=50):
    try:
        return int(float(v))
    except Exception:
        return default

def _append_log(entry: dict, cap: int = 500) -> None:
    pass  # log AI disimpan ke chats; ringkas agar tidak bocor memori

def _fmt_ts(ts):
    if not ts:
        return ""
    try:
        import datetime as _dt
        if isinstance(ts, str):
            parsed = _dt.datetime.fromisoformat(ts.replace("Z", "+00:00"))
        else:
            parsed = ts
        if parsed.tzinfo is None:
            return parsed.strftime("%H:%M")
        return parsed.astimezone().strftime("%H:%M")
    except Exception:
        s = str(ts)
        return s[11:16] if len(s) >= 16 else s

async def _read_json(req: Request) -> dict:
    try:
        d = await req.json()
        return d if isinstance(d, dict) else {}
    except Exception as e:
        raise HTTPException(400, f"Invalid JSON payload: {e}")

# ---------- AUTH ----------
def create_token(user: dict) -> str:
    if not (pyjwt and JWT_SECRET_KEY):
        raise HTTPException(403, "Login disabled: set JWT_SECRET_KEY")
    payload = {
        "sub": user["id"], "role": user["role"], "username": user["username"],
        "name": user["display_name"],
        "exp": datetime.datetime.now(datetime.timezone.utc) + datetime.timedelta(seconds=JWT_EXPIRATION),
    }
    return pyjwt.encode(payload, JWT_SECRET_KEY, algorithm="HS256")

def _dev_user():
    return {"id":"dev","role":"admin","username":"dev","display_name":"Dev","is_active":True,
            "fonnte_token":FONNTE_TOKEN,"fonnte_from_number":FONNTE_FROM,"groq_api_key":GROQ_API_KEY}

async def get_current_user(request: Request, authorization: str = Header(default="")) -> dict:
    """Dependency auth untuk semua /api/v1/*. Dev mode (tanpa key) => user dev admin."""
    if not AUTH_ENABLED:
        return _dev_user()
    scheme, _, token = authorization.partition(" ")
    token = token.strip() if scheme.lower() == "bearer" else authorization.strip()
    if not token and API_TOKEN:
        return _dev_user()
    if not token:
        raise HTTPException(401, "Unauthorized")
    # API_TOKEN static => admin dev
    if API_TOKEN and hmac.compare_digest(token, API_TOKEN):
        return _dev_user()
    if pyjwt and JWT_SECRET_KEY:
        try:
            payload = pyjwt.decode(token, JWT_SECRET_KEY, algorithms=["HS256"])
            u = get_user_by_id(payload.get("sub"))
            if not u:
                raise HTTPException(401, "User tidak ditemukan")
            if not u["is_active"]:
                raise HTTPException(403, "User non-aktif")
            return u
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(401, f"Token invalid: {e}")
    raise HTTPException(401, "Unauthorized")

def require_admin(cu: dict):
    if cu["role"] != "admin":
        raise HTTPException(403, "Hanya admin")

def scope_owner(cu: dict, req_owner: str = None):
    """Admin: req_owner='all' atau kosong => None (semua owner);
       req_owner=<uid> => owner spesifik. User: selalu dirinya."""
    if cu["role"] != "admin":
        return cu["id"]
    if req_owner in (None, "", "all", "me"):
        return None  # semua
    return req_owner

def public_user(u: dict) -> dict:
    return {
        "id": u["id"], "username": u["username"], "role": u["role"],
        "display_name": u["display_name"], "is_active": u["is_active"],
        "created_at": u["created_at"],
        "fonnte_from_number": u.get("fonnte_from_number",""),
        "has_fonnte_token": bool(u.get("fonnte_token")),
        "has_groq_key": bool(u.get("groq_api_key")),
    }

# ---------- RATE LIMITING (in-memory) ----------
_RATE_HITS = {}

def _rate_limited(key: str, limit: int = RATE_MAX, window: int = RATE_WINDOW) -> bool:
    now = time.time()
    hits = [t for t in _RATE_HITS.get(key, []) if now - t < window]
    if len(hits) >= limit:
        _RATE_HITS[key] = hits
        return False
    hits.append(now)
    _RATE_HITS[key] = hits
    return True

# ---------- GROQ CLIENT ----------
try:
    from groq import Groq
except Exception:
    Groq = None

class GroqClient:
    def __init__(self, api_key: str = None):
        key = api_key or GROQ_API_KEY
        self.client = Groq(api_key=key) if (key and Groq) else None
        self.available = self.client is not None
        self.api_key = key

    def _chat_json(self, system: str, user: str, timeout: int = 30) -> dict:
        if not self.available:
            raise RuntimeError("groq not configured")
        resp = self.client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[{"role":"system","content":system + " Output in JSON format."},{"role":"user","content":user}],
            temperature=0.3, max_tokens=800, top_p=0.9,
            response_format={"type":"json_object"},
            timeout=timeout,
        )
        content = resp.choices[0].message.content
        try:
            return json.loads(content)
        except json.JSONDecodeError as e:
            log.warning("Groq JSON decode error: %s | content: %s", e, content[:500])
            raise

    @staticmethod
    def _classify_intent(text: str) -> dict:
        t = text.lower()
        if any(w in t for w in ["harga","biaya","berapa","modal","tarif"]):
            return {"intent":80,"label":"Minta informasi"}
        if any(w in t for w in ["jadwal","kapan","waktu","tanggal"]):
            return {"intent":70,"label":"Tanya jadwal"}
        if any(w in t for w in ["sertifikat","skema","bnsp","asesor"]):
            return {"intent":75,"label":"Tanya sertifikasi"}
        if any(w in t for w in ["promo","diskon","daftar","registrasi","penawaran"]):
            return {"intent":85,"label":"Mau promo/daftar"}
        return {"intent":50,"label":"General question"}

    def analyze(self, message: str, chat_history: str = "") -> dict:
        try:
            def _esc2(s):
                if not s: return ""
                s = str(s)
                return s.replace("{", "{{").replace("}", "}}")
            comps = self._chat_json("Anda adalah AI Sales Assistant profesional. Output in JSON format.",
                                    ASSESSMENT_PROMPT.format(
                                        message=_esc2(message[:2000]),
                                        chat_history=_esc2(chat_history[:2000])))
        except Exception as e:
            log.warning("Groq analyze error -> fallback: %s", e)
            it = self._classify_intent(message)
            comps = {"intent":it["intent"],"product_match":60,"urgency":50,
                     "sentiment":55,"chat_history":50,"decision_maker":70,
                     "intent_label":it["label"],"product":"Umum",
                     "urgency_label":"Sedang","sentiment_label":"Netral",
                     "customer_stage":"Awareness"}
        w = {"intent":.30,"product_match":.20,"urgency":.20,"sentiment":.10,
             "chat_history":.10,"decision_maker":.10}
        score = round(sum(comps.get(k,50)*w[k] for k in w))
        cat  = "Hot Lead" if score>=70 else ("Warm Lead" if score>=40 else "Cold Lead")
        badge= "🟢" if score>=70 else ("🟡" if score>=40 else "🔴")
        return {"lead_score":score,"category":cat,"badge":badge,
                "components":{k:comps.get(k,50) for k in w},
                "intent_label":comps.get("intent_label",""),
                "product":comps.get("product",""),
                "urgency_label":comps.get("urgency_label",""),
                "sentiment_label":comps.get("sentiment_label",""),
                "customer_stage":comps.get("customer_stage",""),
                "analysis_timestamp":datetime.datetime.now(datetime.timezone.utc).isoformat()}

    def generate_reply(self, message: str, context: dict, kb: str = "") -> dict:
        try:
            def _esc(s):
                if not s: return ""
                s = str(s)
                return s.replace("{", "{{").replace("}", "}}")
            hist = context.get("chat_history") or ""
            if isinstance(hist, list):
                hist = "\n".join(str(x) for x in hist[-12:])
            prompt = REPLY_PROMPT.format(
                customer_name=_esc(context.get("customer_name","Customer")),
                message=_esc(message[:2000]),
                chat_history=_esc(hist[:2500] or "(belum ada riwayat)"),
                product=_esc(context.get("product","Umum")),
                stage=_esc(context.get("customer_stage","Awareness")),
                score=context.get("lead_score",50),
                knowledge_chunks=_esc(kb[:3000] or "(tidak ada konteks)")
            )
            out = self._chat_json("Anda adalah AI Sales Assistant profesional. Output in JSON format.", prompt)
            if not isinstance(out, dict):
                raise ValueError(f"Groq returned non-dict: {type(out)}")
            if "suggested_reply" not in out:
                raise KeyError("suggested_reply missing from Groq response")
            return {"suggested_reply":out["suggested_reply"],"confidence_score":min(99,_to_int(out.get("confidence_score",70))),
                    "sources":out.get("sources",[]),"fallback":out.get("fallback","")}
        except Exception as e:
            log.warning("Groq reply error -> fallback: %s | type: %s", e, type(e).__name__)
            it = self._classify_intent(message)
            reply = "Halo Kak 😊\n\nTerima kasih sudah menghubungi ICC.\n\nUntuk memberi informasi yang tepat, boleh Kak beri tahu:"
            if any(w in message.lower() for w in ["harga","biaya"]):
                reply += "  jumlah peserta, kebutuhan sertifikasi, dan lokasi perusahaan?"
            elif any(w in message.lower() for w in ["jadwal","kapan"]):
                reply += "  lokasi / preferensi wilayah?"
            else:
                reply += "  produk apa yang Kakak butuhkan?"
            return {"suggested_reply":reply,"confidence_score":min(99,_to_int(it["intent"],50)+10),
                    "sources":[{"type":"faq","reference":f"FAQ - {context.get('product','Umum')}"}],
                    "fallback":"Maaf kak, saya kurang yakin. Sales kami akan menghubungi segera."}

def groq_for(user: dict) -> GroqClient:
    return GroqClient(api_key=(user or {}).get("groq_api_key") or GROQ_API_KEY)

# ---------- PROMPTS ----------
ASSESSMENT_PROMPT = """Anda adalah AI Sales Assistant profesional untuk perusahaan training & sertifikasi (ICC Holding).
Tugas: analisis pesan customer → skor 6 komponen (0-100 tiap) + label.
Komponen:
1. Intent Signal (niat beli/tanya harga/daftar)
2. Product Match (kesesuaian produk yg ditawarkan)
3. Timeline Urgency (urgen/bulan ini/belum tentu)
4. Engagement Level (balas cepat/panjang/pendek)
5. Chat History Depth (baru/ada riwayat)
6. Decision Maker Authority (pemutus/ pengaruh/ staff)

Output HANYA JSON (tanpa teks lain):
{{"intent": 0-100, "product_match": 0-100, "urgency": 0-100, "sentiment": 0-100,
"chat_history": 0-100, "decision_maker": 0-100, "intent_label": "...",
"product": "...", "urgency_label": "...", "sentiment_label": "...",
"customer_stage": "..."}}

Pesan customer: {message}
Riwayat chat: {chat_history}"""

REPLY_PROMPT = """Anda Sales Assistant senior ICC Holding (LSP BNSP terakreditasi). Produk: POPAL, GIS/Geomatika, Sertifikasi BNSP, Pelatihan K3 (AK3, H2S, Confined Space, Working at Height), ISO, Manajemen Proyek.

Customer: {customer_name}
Pesan terbaru: {message}
Riwayat chat:
{chat_history}
Produk: {product}
Stage: {stage}
Score: {score}

ATURAN:
1. Jawab berdasarkan riwayat chat + pesan terbaru (lanjutkan percakapan, JANGAN ulangi sapaan/pertanyaan yang sudah dijawab)
2. Jawab spesifik + tanyakan data yang belum ada: jumlah peserta, lokasi, timeline, skema
3. Harga: range/estimasi saja (jangan final tanpa konteks)
4. Cross-sell cerdas: POPAL+TOT, K3+Sertif BNSP K3, GIS+Surveyor, Corporate Package
5. Value prop: LSP BNSP langsung, sertifikat resmi database nasional, compliance UU/PP/Kemenaker
6. Tone: konsultatif partner solusi, emoji max 2, bahasa Indonesia

Knowledge: {knowledge_chunks}

Output HANYA JSON:
{{"suggested_reply":"...", "confidence_score":0-100, "sources":[{{"type":"faq|sop|product|regulation","reference":"..."}}], "fallback":"..."}}
"""

# ---------- WebSocket ----------
class ConnectionManager:
    def __init__(self): self.connections = []
    async def connect(self, ws): await ws.accept(); self.connections.append(ws)
    def disconnect(self, ws):
        if ws in self.connections: self.connections.remove(ws)
    async def broadcast(self, data: dict):
        dead = []
        for ws in self.connections:
            try: await ws.send_json(data)
            except Exception: dead.append(ws)
        for ws in dead: self.disconnect(ws)

manager = ConnectionManager()

@app.websocket("/ws")
async def ws_endpoint(ws: WebSocket):
    await manager.connect(ws)
    try:
        while True: await ws.receive_text()
    except WebSocketDisconnect:
        manager.disconnect(ws)

# ---------- FONNTE ----------
def fonnte_token_for(owner: dict) -> str:
    return (owner or {}).get("fonnte_token") or FONNTE_TOKEN

async def fonnte_send(target: str, text: str, token: str = None) -> dict:
    """Kirim pesan WA via Fonnte API. Token diambil dari user owner (fallback env)."""
    tk = token or FONNTE_TOKEN
    if not tk:
        return {"status_code":0,"error":"FONNTE_TOKEN not set"}
    async with httpx.AsyncClient(timeout=15) as client:
        r = await client.post("https://api.fonnte.com/send",
            data={"target":target,"message":text,"countryCode":"62"},
            headers={"Authorization":tk})
        return {"status_code":r.status_code,"response":r.text}

# ---------- KNOWLEDGE BASE ----------
def chunk_text(text: str, size: int = KB_CHUNK_SIZE, overlap: int = KB_CHUNK_OVERLAP) -> list:
    text = (text or "").strip()
    if not text:
        return []
    parts, buf = [], []
    for para in text.replace("\r\n", "\n").split("\n\n"):
        para = para.strip()
        if not para:
            continue
        if len(para) <= size:
            parts.append(para)
        else:
            start = 0
            while start < len(para):
                end = min(start + size, len(para))
                parts.append(para[start:end].strip())
                if end >= len(para):
                    break
                start = max(end - overlap, start + 1)
    out = []
    for p in parts:
        if out and len(out[-1]) + 1 + len(p) <= size:
            out[-1] = out[-1] + "\n" + p
        else:
            out.append(p)
    return [{"idx": i, "text": c, "chars": len(c)} for i, c in enumerate(out) if c]

def _tokens(q: str) -> set:
    stop = {"yang","dan","atau","untuk","dari","dengan","pada","ini","itu","ada","saya","kami","kak","apakah","berapa","mohon","info","tentang","the","a","an","of","to","in"}
    return {w for w in "".join(ch.lower() if ch.isalnum() else " " for ch in (q or "")).split() if len(w) >= 3 and w not in stop}

def score_chunk(query: str, chunk: str, doc_name: str = "", category: str = "") -> float:
    toks = _tokens(query)
    if not toks:
        return 0.0
    blob = f"{doc_name} {category} {chunk}".lower()
    hits = sum(1 for t in toks if t in blob)
    score = hits / max(len(toks), 1)
    words = [w for w in "".join(ch.lower() if ch.isalnum() else " " for ch in query).split() if len(w) >= 3]
    for i in range(len(words) - 1):
        phrase = f"{words[i]} {words[i+1]}"
        if phrase in blob:
            score += 0.25
    if category and category.lower() in (query or "").lower():
        score += 0.15
    return round(score, 4)

def _kb_rows_for_viewer(cu: dict, owner_filter: str = None) -> list:
    """Dokumen KB yang boleh dilih oleh cu. User: miliknya + global ('*').
       Admin all: semua. Admin spesifik: milik user itu + global."""
    if cu["role"] == "admin" and owner_filter in (None, "all"):
        return db_query("SELECT * FROM knowledge_base ORDER BY uploaded_at DESC LIMIT 500")
    if cu["role"] == "admin" and owner_filter:
        return db_query("SELECT * FROM knowledge_base WHERE owner_id=? OR owner_id='*' ORDER BY uploaded_at DESC LIMIT 500", (owner_filter,))
    # user
    return db_query("SELECT * FROM knowledge_base WHERE owner_id=? OR owner_id='*' ORDER BY uploaded_at DESC LIMIT 500", (cu["id"],))

def search_kb_chunks(query: str, category: str = None, limit: int = 5, cu: dict = None, owner_filter: str = None) -> list:
    rows = _kb_rows_for_viewer(cu or _dev_user(), owner_filter)
    if category and category in PRODUCT_CATEGORIES:
        rows = [r for r in rows if r.get("category") == category]
    scored = []
    for row in rows:
        chunks = row.get("chunks") or []
        if isinstance(chunks, str):
            try: chunks = json.loads(chunks)
            except Exception: chunks = []
        if not chunks and row.get("kb_text"):
            chunks = chunk_text(row["kb_text"])
        for ch in chunks:
            text = ch.get("text") if isinstance(ch, dict) else str(ch)
            if not text:
                continue
            s = score_chunk(query, text, row.get("name",""), row.get("category",""))
            if s <= 0:
                continue
            scored.append({
                "score": s, "text": text[:1200], "name": row.get("name",""),
                "category": row.get("category",""), "doc_id": row.get("id"),
                "chunk_idx": ch.get("idx") if isinstance(ch, dict) else None,
                "file_url": row.get("file_url",""),
            })
    scored.sort(key=lambda x: x["score"], reverse=True)
    out, seen = [], set()
    for item in scored:
        key = item["text"][:120]
        if key in seen:
            continue
        seen.add(key)
        out.append(item)
        if len(out) >= limit:
            break
    return out

def format_kb_context(chunks: list) -> str:
    if not chunks:
        return ""
    parts = []
    for i, c in enumerate(chunks, 1):
        parts.append(f"[{i}] ({c.get('category','?')} | {c.get('name','?')} | score={c.get('score',0)})\n{c.get('text','')}")
    return "\n\n".join(parts)

async def upload_to_storage(file: UploadFile) -> dict:
    """Upload file produk ke Supabase Storage (jika ada), else mock."""
    if not HAS_SUPABASE:
        content = await file.read()
        return {"status":"mock","filename":file.filename, "size": len(content)}
    try:
        content = await file.read()
        ext = os.path.splitext(file.filename or "")[1]
        fname = f"{uuid.uuid4().hex}{ext}"
        res = supabase_admin.storage.from_("products").upload(fname, content)
        if res.status_code == 200:
            public_url = supabase_admin.storage.from_("products").get_public_url(fname)
            return {"status":"uploaded","filename":fname,"public_url":public_url}
        return {"status":"error","code":res.status_code,"response":getattr(res, 'text', str(res))}
    except Exception as e:
        return {"status":"error","error":str(e)}

# ---------- CUSTOMER / CHAT STORE (SQLite) ----------
def ensure_customer(owner_id, phone, name=None, analysis=None):
    now = datetime.datetime.now(datetime.timezone.utc).isoformat()
    existing = db_query("SELECT * FROM customers WHERE owner_id=? AND phone=?", (owner_id, phone))
    if existing:
        if name and existing[0]["name"] in (None, "", phone):
            db_exec("UPDATE customers SET name=?, last_updated=? WHERE id=?", (name, now, existing[0]["id"]))
        if analysis:
            db_exec("UPDATE customers SET last_score=?, category=?, badge=?, product=?, last_updated=? WHERE id=?",
                    (analysis.get("lead_score",0), analysis.get("category","Cold Lead"), analysis.get("badge","🔴"),
                     analysis.get("product","Umum"), now, existing[0]["id"]))
        return db_query("SELECT * FROM customers WHERE owner_id=? AND phone=?", (owner_id, phone))[0]
    db_exec("INSERT INTO customers (owner_id,phone,name,last_score,category,badge,product,created_at,last_updated) "
            "VALUES (?,?,?,?,?,?,?,?,?)",
            (owner_id, phone, name or phone,
             (analysis or {}).get("lead_score",0), (analysis or {}).get("category","Cold Lead"),
             (analysis or {}).get("badge","🔴"), (analysis or {}).get("product","Umum"), now, now))
    return db_query("SELECT * FROM customers WHERE owner_id=? AND phone=?", (owner_id, phone))[0]

def append_chat(owner_id, phone, text, direction, analysis=None, name=None):
    now = datetime.datetime.now(datetime.timezone.utc).isoformat()
    db_exec("INSERT INTO chats (owner_id,phone,text,direction,timestamp,lead_score,category,intent_label,sender_name) "
            "VALUES (?,?,?,?,?,?,?,?,?)",
            (owner_id, phone, text, direction, now,
             (analysis or {}).get("lead_score"), (analysis or {}).get("category"),
             (analysis or {}).get("intent_label"), name))
    if direction == "in":
        db_exec("UPDATE customers SET unread=unread+1, last_message=?, last_updated=? WHERE owner_id=? AND phone=?",
                (text[:200], now, owner_id, phone))
    else:
        db_exec("UPDATE customers SET unread=0, last_message=?, last_updated=? WHERE owner_id=? AND phone=?",
                (text[:200], now, owner_id, phone))

def get_customer_dict(row, msgs=None):
    return {
        "phone": row["phone"], "name": row.get("name") or row["phone"],
        "owner_id": row.get("owner_id"),
        "score": row.get("last_score",0) or 0, "category": row.get("category") or "Cold Lead",
        "badge": row.get("badge") or "🔴", "product": row.get("product") or "Umum",
        "last": row.get("last_message") or "", "unread": row.get("unread",0) or 0,
        "created": row.get("created_at"), "stage": "Awareness", "msgs": msgs or [],
    }

def process_incoming(payload: dict, owner: dict):
    sender = payload.get("sender") or payload.get("pengirim") or payload.get("from")
    text   = payload.get("message") or payload.get("pesan") or payload.get("text")
    name   = payload.get("name") or payload.get("pushname") or sender
    if not sender or not text:
        return None, None, None
    g = groq_for(owner)
    analysis = g.analyze(text)
    cust_row = ensure_customer(owner["id"], sender, name, analysis)
    append_chat(owner["id"], sender, text, "in", analysis, name)
    cust = get_customer_dict(cust_row)
    msg = {"from":sender,"text":text,"timestamp":datetime.datetime.now(datetime.timezone.utc).isoformat()}
    log.info("[incoming] owner=%s %s (%s) -> %s (score=%s)", owner["id"], sender, name, text[:60], analysis["lead_score"])
    return cust, msg, analysis

# ---------- ROUTES ----------
@app.get("/", include_in_schema=False)
async def serve_frontend():
    return FileResponse(FRONTEND_FILE)

@app.get("/health")
async def health():
    users = db_query("SELECT COUNT(*) c FROM users")
    return {"status":"ok","version":"0.5.0","multi_user":True,
            "groq_ready":bool(GROQ_API_KEY or any(u["groq_api_key"] for u in list_users())),
            "groq_model":GROQ_MODEL,
            "fonnte_ready":bool(FONNTE_TOKEN or any(u["fonnte_token"] for u in list_users())),
            "db":"sqlite","user_count":users[0]["c"] if users else 0,
            "auth_enabled":AUTH_ENABLED,
            "webhook_secured":bool(WEBHOOK_SECRET),
            "rate_max":RATE_MAX,"rate_window":RATE_WINDOW}

# ---------- AUTH ----------
@app.post("/api/v1/auth/login")
async def auth_login(req: Request):
    if not (pyjwt and JWT_SECRET_KEY):
        raise HTTPException(403, "Login disabled: set JWT_SECRET_KEY")
    d = await _read_json(req)
    u, p = (d.get("username") or "").strip(), (d.get("password") or "").strip()
    logging.info("LOGIN_ATTEMPT username=%r pwlen=%d", u, len(p))
    user = get_user_by_username(u)
    if not user or not verify_password(p, _pw_hash_for(u)):
        raise HTTPException(401, "Invalid credentials")
    if not user["is_active"]:
        raise HTTPException(403, "User non-aktif")
    token = create_token(user)
    return {"status":"success","token":token,"expires_in":JWT_EXPIRATION,
            "user":{"id":user["id"],"username":user["username"],"role":user["role"],"display_name":user["display_name"]}}

def _pw_hash_for(username):
    rows = db_query("SELECT password_hash FROM users WHERE username=?", (username,))
    return rows[0]["password_hash"] if rows else ""

@app.get("/api/v1/me")
async def me(cu: dict = Depends(get_current_user)):
    return {"status":"success","data":public_user(cu)}

@app.put("/api/v1/me")
async def update_me(req: Request, cu: dict = Depends(get_current_user)):
    d = await _read_json(req)
    fields = {}
    if "display_name" in d: fields["display_name"] = str(d["display_name"])[:80]
    if "fonnte_token" in d: fields["fonnte_token"] = str(d["fonnte_token"])
    if "fonnte_from_number" in d: fields["fonnte_from_number"] = str(d["fonnte_from_number"])
    if d.get("new_password"):
        if not d.get("current_password") or not verify_password(d["current_password"], _pw_hash_for(cu["username"])):
            raise HTTPException(400, "Password saat ini salah")
        set_password(cu["id"], str(d["new_password"]))
    if fields:
        update_user(cu["id"], **fields)
    return {"status":"success","data":public_user(get_user_by_id(cu["id"]))}

# ---------- ADMIN: USER MANAGEMENT ----------
@app.get("/api/v1/admin/users")
async def admin_list_users(cu: dict = Depends(get_current_user)):
    require_admin(cu)
    return {"status":"success","data":[public_user(u) for u in list_users()]}

@app.post("/api/v1/admin/users")
async def admin_create_user(req: Request, cu: dict = Depends(get_current_user)):
    require_admin(cu)
    d = await _read_json(req)
    uname = (d.get("username") or "").strip()
    pw = d.get("password") or ""
    role = "admin" if d.get("role") == "admin" else "user"
    if not uname or not pw:
        raise HTTPException(400, "username & password wajib")
    if get_user_by_username(uname):
        raise HTTPException(409, "Username sudah dipakai")
    u = create_user(uname, pw, role,
                    display_name=d.get("display_name") or uname,
                    fonnte_token=d.get("fonnte_token") or "",
                    fonnte_from=d.get("fonnte_from_number") or "")
    return {"status":"success","data":public_user(u)}

@app.put("/api/v1/admin/users/{uid}")
async def admin_update_user(uid: str, req: Request, cu: dict = Depends(get_current_user)):
    require_admin(cu)
    if uid == cu["id"] and (req and (await _read_json(req)).get("role") == "user"):
        # cegah admin menurunkan dirinya sendiri
        raise HTTPException(400, "Tidak bisa menurunkan role sendiri")
    d = await _read_json(req)
    fields = {}
    if "role" in d: fields["role"] = "admin" if d["role"] == "admin" else "user"
    if "display_name" in d: fields["display_name"] = str(d["display_name"])[:80]
    if "fonnte_token" in d: fields["fonnte_token"] = str(d["fonnte_token"])
    if "fonnte_from_number" in d: fields["fonnte_from_number"] = str(d["fonnte_from_number"])
    if "is_active" in d: fields["is_active"] = 1 if d["is_active"] else 0
    if d.get("new_password"):
        set_password(uid, str(d["new_password"]))
    if fields:
        update_user(uid, **fields)
    u = get_user_by_id(uid)
    if not u:
        raise HTTPException(404, "User tidak ditemukan")
    return {"status":"success","data":public_user(u)}

@app.delete("/api/v1/admin/users/{uid}")
async def admin_delete_user(uid: str, cu: dict = Depends(get_current_user)):
    require_admin(cu)
    if uid == cu["id"]:
        raise HTTPException(400, "Tidak bisa menghapus diri sendiri")
    if not get_user_by_id(uid):
        raise HTTPException(404, "User tidak ditemukan")
    delete_user(uid)
    return {"status":"success","deleted":uid}

# ---------- WHATSAPP WEBHOOK ----------
def _webhook_owner(req: Request) -> dict:
    uid = req.query_params.get("uid")
    if uid:
        u = get_user_by_id(uid)
        if u and u["is_active"]:
            return u
    # fallback dev: user pertama
    rows = db_query("SELECT * FROM users LIMIT 1")
    return _row_to_user(rows[0]) if rows else _dev_user()

def _webhook_ok(request: Request) -> bool:
    if not WEBHOOK_SECRET:
        return True
    header = request.headers.get("x-webhook-token", "")
    query  = request.query_params.get("token", "")
    try:
        if header and hmac.compare_digest(header, WEBHOOK_SECRET):
            return True
        if query and hmac.compare_digest(query, WEBHOOK_SECRET):
            return True
    except Exception:
        return False
    return False

@app.post("/webhook/fonnte")
async def webhook_fonnte(req: Request):
    if not _webhook_ok(req):
        raise HTTPException(401, "Invalid webhook token")
    payload = await _read_json(req)
    client_ip = req.client.host if req.client else "?"
    if not _rate_limited(f"wh:{client_ip}:{payload.get('sender','')}", limit=120):
        raise HTTPException(429, "Too many webhook requests")
    if "stateid" in payload or ("device" in payload and "message" not in payload and "sender" not in payload):
        return {"status":"received","type":"device_status"}
    owner = _webhook_owner(req)
    cust, msg, analysis = process_incoming(payload, owner)
    if cust is None:
        return {"status":"received","type":"skip"}
    await manager.broadcast({"type":"chat_incoming","customer":cust,"message":msg,"analysis":analysis,
                             "owner_id":owner["id"]})
    return {"status":"received","type":"message"}

@app.post("/webhook/fonte")
async def webhook_fonte(req: Request):
    return await webhook_fonnte(req)

# ---------- ANALYTICS / REPLY ----------
@app.post("/api/v1/assessment/analyze")
async def assess(req: Request, cu: dict = Depends(get_current_user)):
    d = await _read_json(req)
    g = groq_for(cu)
    analysis = g.analyze(d.get("message",""), d.get("chat_history",""))
    return {"status":"success","data":analysis}

@app.post("/api/v1/reply/generate")
async def reply_generate(req: Request, cu: dict = Depends(get_current_user)):
    d = await _read_json(req)
    kb = d.get("knowledge_chunks","")
    product = d.get("context",{}).get("product","Umum")
    message = d.get("message","")
    hits = []
    if not kb:
        cat_match = None
        for cat in PRODUCT_CATEGORIES:
            if cat.lower() in (product or "").lower() or cat.lower() in (message or "").lower():
                cat_match = cat
                break
        owner_filter = scope_owner(cu, d.get("owner"))
        hits = search_kb_chunks(message or product, category=cat_match, limit=5, cu=cu, owner_filter=owner_filter)
        if not hits and cat_match:
            hits = search_kb_chunks(message or product, category=None, limit=5, cu=cu, owner_filter=owner_filter)
        kb = format_kb_context(hits)
    g = groq_for(cu)
    reply = g.generate_reply(message, d.get("context",{}), kb)
    if isinstance(reply, dict) and "suggested_reply" in reply:
        reply["kb_hits"] = [{"name":h.get("name"),"category":h.get("category"),"score":h.get("score"),"chunk_idx":h.get("chunk_idx")} for h in hits]
        extra = [{"type":"kb","reference":f"{h.get('category','')}/{h.get('name','')}#{h.get('chunk_idx')}"} for h in hits[:3]]
        reply["sources"] = (reply.get("sources") or []) + extra
        reply["confidence"] = reply.get("confidence_score", 0)
    else:
        reply = {"suggested_reply":"(fallback error)","confidence_score":20,"sources":[],"fallback":"error","kb_hits":[]}
    return {"status":"success","data":reply}

# ---------- MESSAGES ----------
@app.post("/api/v1/messages/send")
async def send_message(req: Request, cu: dict = Depends(get_current_user)):
    d = await _read_json(req)
    to, text = d.get("to",""), d.get("text","")
    if not to or not text:
        raise HTTPException(400, "to & text wajib")
    # Tentukan owner token: admin bisa kirim atas nama user lain (owner di body)
    owner_id = scope_owner(cu, d.get("owner"))
    if owner_id is None:
        owner_id = cu["id"]  # admin tanpa owner spesifik -> kirim pakai token sendiri
    owner = cu if owner_id == cu["id"] else (get_user_by_id(owner_id) or cu)
    token = fonnte_token_for(owner)
    client_ip = req.client.host if req.client else "?"
    if not _rate_limited(f"send:{client_ip}:{to}"):
        raise HTTPException(429, "Too many messages, please slow down")
    res = await fonnte_send(to, text, token)
    append_chat(owner["id"], to, text, "out")
    return {"status":"success","sent":res["status_code"]==200,"fonnte":res,"owner_id":owner["id"]}

# ---------- CUSTOMERS ----------
@app.get("/api/v1/customers/{phone}/messages")
async def customer_messages(phone: str, req: Request, cu: dict = Depends(get_current_user)):
    owner_id = scope_owner(cu, req.query_params.get("owner"))
    if owner_id is None:
        rows = db_query("SELECT * FROM chats WHERE phone=? ORDER BY timestamp ASC LIMIT 200", (phone,))
    else:
        rows = db_query("SELECT * FROM chats WHERE owner_id=? AND phone=? ORDER BY timestamp ASC LIMIT 200", (owner_id, phone))
    msgs = [{"dir":(r.get("direction") or "in"),"t":_fmt_ts(r.get("timestamp")),"d":(r.get("text") or ""),"ts":r.get("timestamp")} for r in rows]
    if owner_id is not None:
        db_exec("UPDATE customers SET unread=0 WHERE owner_id=? AND phone=?", (owner_id, phone))
    else:
        db_exec("UPDATE customers SET unread=0 WHERE phone=?", (phone,))
    return {"status":"success","data":msgs,"count":len(msgs)}

@app.get("/api/v1/customers")
async def list_customers(req: Request, cu: dict = Depends(get_current_user)):
    owner_id = scope_owner(cu, req.query_params.get("owner"))
    if owner_id is None:
        rows = db_query("SELECT * FROM customers ORDER BY last_updated DESC LIMIT 500")
    else:
        rows = db_query("SELECT * FROM customers WHERE owner_id=? ORDER BY last_updated DESC LIMIT 500", (owner_id,))
    return [get_customer_dict(r) for r in rows]

# ---------- STATS ----------
@app.get("/api/v1/stats")
async def get_stats(req: Request, cu: dict = Depends(get_current_user)):
    owner_id = scope_owner(cu, req.query_params.get("owner"))
    now = datetime.datetime.now(datetime.timezone.utc)
    hot = warm = cold = 0
    chat_in = 0
    product_counts = {}
    activity = {}
    if owner_id is None:
        cust_rows = db_query("SELECT category,product FROM customers")
        chat_rows = db_query("SELECT timestamp FROM chats WHERE direction='in'")
    else:
        cust_rows = db_query("SELECT category,product FROM customers WHERE owner_id=?", (owner_id,))
        chat_rows = db_query("SELECT timestamp FROM chats WHERE owner_id=? AND direction='in'", (owner_id,))
    for r in cust_rows:
        cat = r.get("category") or "Cold Lead"
        if "Hot" in cat: hot += 1
        elif "Warm" in cat: warm += 1
        else: cold += 1
        prod = r.get("product") or "Umum"
        product_counts[prod] = product_counts.get(prod, 0) + 1
    for r in chat_rows:
        ts = r.get("timestamp")
        if ts:
            day = ts[:10]
            activity[day] = activity.get(day, 0) + 1
            chat_in += 1
    total = hot + warm + cold
    days = []
    for i in range(6, -1, -1):
        ddate = (now - datetime.timedelta(days=i)).date()
        label = ["Sen", "Sel", "Rab", "Kam", "Jum", "Sab", "Min"][ddate.weekday()]
        days.append({"label": label, "date": ddate.isoformat(), "value": activity.get(ddate.isoformat(), 0)})
    top_products = [{"name": k, "count": v, "sub": ""} for k, v in
                    sorted(product_counts.items(), key=lambda x: x[1], reverse=True)[:5]]
    return {"status": "success", "data": {
        "owner_id": owner_id, "scope": "all" if owner_id is None else owner_id,
        "chat_count": chat_in,
        "hot": hot, "warm": warm, "cold": cold, "total": total,
        "distribution": {
            "hot": round(hot / total * 100) if total else 0,
            "warm": round(warm / total * 100) if total else 0,
            "cold": round(cold / total * 100) if total else 0,
        },
        "avg_response_sec": None,
        "activity": days,
        "top_products": top_products,
        "response_trend": [],
        "generated_at": now.isoformat(),
    }}

# ---------- PRODUCTS ----------
@app.post("/api/v1/products")
async def create_product(p: dict, cu: dict = Depends(get_current_user)):
    pid = p.get("id") or str(uuid.uuid4())
    data = {"id":pid,"owner_id":cu["id"],"name":p["name"],"category":p.get("category","Umum"),
            "description":p.get("description",""),"price_range":p.get("price_range",""),
            "duration":p.get("duration",""),"kb_text":p.get("kb_text",""),
            "created_at":datetime.datetime.now(datetime.timezone.utc).isoformat()}
    db_exec("INSERT OR REPLACE INTO products (id,owner_id,name,category,description,price_range,duration,kb_text,created_at) "
            "VALUES (:id,:owner_id,:name,:category,:description,:price_range,:duration,:kb_text,:created_at)", data)
    return {"status":"saved","product_id":pid}

@app.get("/api/v1/products")
async def list_products(req: Request, category: str = None, cu: dict = Depends(get_current_user)):
    owner_id = scope_owner(cu, req.query_params.get("owner"))
    if owner_id is None:
        rows = db_query("SELECT * FROM products ORDER BY created_at DESC")
    else:
        rows = db_query("SELECT * FROM products WHERE owner_id=? ORDER BY created_at DESC", (owner_id,))
    if category:
        rows = [r for r in rows if r.get("category") == category]
    return [dict(r) for r in rows]

@app.get("/api/v1/products/{pid}")
async def get_product(pid: str, cu: dict = Depends(get_current_user)):
    rows = db_query("SELECT * FROM products WHERE id=?", (pid,))
    if not rows:
        return {"status":"error","id":pid}
    return dict(rows[0])

# ---------- UPLOAD / KNOWLEDGE ----------
@app.post("/api/v1/upload")
async def upload_product_file(file: UploadFile = File(...), cu: dict = Depends(get_current_user)):
    result = await upload_to_storage(file)
    return result

def _extract_text(content_bytes, ext, file):
    kb_text = ""
    ext_lower = ext
    try:
        if ext_lower == ".txt" or ext_lower in [".md"]:
            kb_text = content_bytes.decode("utf-8", errors="replace")
        elif ext_lower == ".pdf":
            try:
                from pypdf import PdfReader
                reader = PdfReader(io.BytesIO(content_bytes))
                kb_text = "\n\n".join((page.extract_text() or "") for page in reader.pages)
            except ImportError:
                import fitz
                doc = fitz.open(stream=content_bytes, filetype="pdf")
                kb_text = "\n\n".join(page.get_text() for page in doc)
                doc.close()
        elif ext_lower == ".docx":
            from docx import Document
            doc = Document(io.BytesIO(content_bytes))
            kb_text = "\n\n".join(p.text or "" for p in doc.paragraphs if p.text.strip())
        elif ext_lower in [".xlsx", ".xls"]:
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(content_bytes), read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        rows.append("\t".join(cells))
            kb_text = "\n".join(rows)
            wb.close()
        elif ext_lower in [".pptx"]:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content_bytes))
            kb_text = "\n\n".join(shape.text or "" for slide in prs.slides for shape in slide.shapes if hasattr(shape,"text"))
    except ImportError as ie:
        log.warning("[KB Upload] Missing dep for %s: %s", ext, ie)
    except Exception as e:
        log.warning("[KB Upload] %s extraction error: %s", ext, e)
        kb_text = f"(extraction error: {e})"
    if not kb_text:
        kb_text = f"File KB: {file.filename}. Buka dokumen untuk detail."
    return kb_text

@app.post("/api/v1/knowledge/upload")
async def upload_knowledge(file: UploadFile = File(...), category: str = Form("Umum"), name: str = Form(""),
                           extract_text: bool = Form(True), cu: dict = Depends(get_current_user)):
    if category not in PRODUCT_CATEGORIES:
        category = "Umum"
    content_bytes = await file.read()
    ext = os.path.splitext(file.filename or "")[1].lower()
    kb_text = _extract_text(content_bytes, ext, file) if extract_text else f"File KB: {file.filename}"
    # storage (optional)
    try:
        reuse_file = UploadFile(filename=file.filename, file=io.BytesIO(content_bytes), headers={})
        upload_result = await upload_to_storage(reuse_file)
    except Exception as e:
        upload_result = {"status":"error","error":str(e)}
    chunks = chunk_text(kb_text)
    doc_id = str(uuid.uuid4())
    now = datetime.datetime.now(datetime.timezone.utc).isoformat()
    db_exec("INSERT INTO knowledge_base (id,owner_id,category,name,filename,file_url,kb_text,chunks,chunk_count,uploaded_at) "
            "VALUES (?,?,?,?,?,?,?,?,?,?)",
            (doc_id, cu["id"], category, name or file.filename or "", file.filename or "",
             upload_result.get("public_url",""), kb_text[:50000], json.dumps(chunks), len(chunks), now))
    return {"status":"saved","doc_id":doc_id,"kb_text_length":len(kb_text),"chunk_count":len(chunks),
            "chunks_preview":chunks[:3],"file":upload_result}

@app.get("/api/v1/knowledge/search")
async def knowledge_search(req: Request, q: str = "", cat: str = None, limit: int = 5, cu: dict = Depends(get_current_user)):
    if not q.strip():
        return {"status":"error","error":"q wajib"}
    category = None
    if cat and cat.strip() and cat.strip() in PRODUCT_CATEGORIES:
        category = cat.strip()
    owner_filter = scope_owner(cu, req.query_params.get("owner"))
    hits = search_kb_chunks(q, category=category, limit=max(1, min(limit, 20)), cu=cu, owner_filter=owner_filter)
    return {"status":"success","data":{"query":q,"category":category,"count":len(hits),"hits":hits}}

def _kb_doc_dict(p: dict) -> dict:
    ch = p.get("chunks") or []
    if isinstance(ch, str):
        try: ch = json.loads(ch)
        except Exception: ch = []
    return {"id":p.get("id"),"owner_id":p.get("owner_id"),"name":p.get("name"),"category":p.get("category"),
           "filename":p.get("filename",""),"url":p.get("file_url"),"uploaded_at":p.get("uploaded_at"),
           "chunk_count":p.get("chunk_count") or len(ch),"preview":(p.get("kb_text") or "")[:400],
           "kb_text":p.get("kb_text") or ""}

@app.get("/api/v1/knowledge")
async def list_knowledge(req: Request, category: str = None, cu: dict = Depends(get_current_user)):
    if category and category not in PRODUCT_CATEGORIES:
        return {"status":"error","error":f"Invalid category. Valid: {PRODUCT_CATEGORIES}"}
    rows = _kb_rows_for_viewer(cu, scope_owner(cu, req.query_params.get("owner")))
    if category:
        rows = [r for r in rows if r.get("category") == category]
    docs = [_kb_doc_dict(r) for r in rows]
    counts = {c:0 for c in PRODUCT_CATEGORIES}
    for d in docs:
        counts[d.get("category")] = counts.get(d.get("category"),0) + 1
    total_chunks = sum(d.get("chunk_count") or 0 for d in docs)
    return {"status":"success","data":{"category":category,"count":len(docs),"docs":docs,
            "category_counts":counts,"total_chunks":total_chunks}}

@app.get("/api/v1/knowledge/doc/{doc_id}")
async def get_knowledge_doc(doc_id: str, cu: dict = Depends(get_current_user)):
    rows = db_query("SELECT * FROM knowledge_base WHERE id=?", (doc_id,))
    if not rows:
        return {"status":"error","error":"doc not found"}
    r = rows[0]
    if cu["role"] != "admin" and r["owner_id"] not in (cu["id"], "*"):
        raise HTTPException(403, "Tidak punya akses")
    return {"status":"success","data":_kb_doc_dict(r)}

@app.put("/api/v1/knowledge/{doc_id}")
async def update_knowledge(doc_id: str, p: dict, cu: dict = Depends(get_current_user)):
    rows = db_query("SELECT * FROM knowledge_base WHERE id=?", (doc_id,))
    if not rows:
        return {"status":"error","error":"doc not found"}
    r = rows[0]
    if cu["role"] != "admin" and r["owner_id"] not in (cu["id"], "*"):
        raise HTTPException(403, "Tidak punya akses")
    name = str(p.get("name","")).strip()
    category = str(p.get("category","")).strip()
    kb_text = str(p.get("kb_text","") or "").strip()
    if category and category not in PRODUCT_CATEGORIES:
        return {"status":"error","error":f"Invalid category. Valid: {PRODUCT_CATEGORIES}"}
    if not name and not kb_text:
        return {"status":"error","error":"name atau kb_text wajib"}
    chunks = chunk_text(kb_text) if kb_text else None
    sets, params = [], []
    if name: sets.append("name=?"); params.append(name)
    if category: sets.append("category=?"); params.append(category)
    if kb_text:
        sets.append("kb_text=?"); params.append(kb_text[:50000])
        sets.append("chunks=?"); params.append(json.dumps(chunks))
        sets.append("chunk_count=?"); params.append(len(chunks))
    params.append(doc_id)
    db_exec(f"UPDATE knowledge_base SET {','.join(sets)} WHERE id=?", tuple(params))
    return {"status":"saved","id":doc_id}

@app.delete("/api/v1/knowledge/{doc_id}")
async def delete_knowledge(doc_id: str, cu: dict = Depends(get_current_user)):
    rows = db_query("SELECT * FROM knowledge_base WHERE id=?", (doc_id,))
    if not rows:
        return {"status":"error","error":"doc not found"}
    r = rows[0]
    if cu["role"] != "admin" and r["owner_id"] not in (cu["id"], "*"):
        raise HTTPException(403, "Tidak punya akses")
    db_exec("DELETE FROM knowledge_base WHERE id=?", (doc_id,))
    return {"status":"deleted","id":doc_id}

@app.get("/api/v1/knowledge/{category}")
async def get_knowledge(category: str, req: Request, cu: dict = Depends(get_current_user)):
    if category not in PRODUCT_CATEGORIES:
        return {"status":"error","error":f"Invalid category. Valid: {PRODUCT_CATEGORIES}"}
    rows = _kb_rows_for_viewer(cu, scope_owner(cu, req.query_params.get("owner")))
    rows = [r for r in rows if r.get("category") == category]
    docs = []
    for p in rows:
        ch = p.get("chunks") or []
        if isinstance(ch, str):
            try: ch = json.loads(ch)
            except Exception: ch = []
        docs.append({"id":p.get("id"),"name":p.get("name"),"category":p.get("category"),
                    "url":p.get("file_url"),"uploaded_at":p.get("uploaded_at"),
                    "chunk_count":p.get("chunk_count") or len(ch),"preview":(p.get("kb_text") or "")[:400]})
    return {"category":category,"docs":docs,"count":len(docs)}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
