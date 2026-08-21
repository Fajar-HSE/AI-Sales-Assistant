import datetime
import io
import json
import logging
import os
import uuid

from fastapi import APIRouter, Depends, File, Form, HTTPException, Request, UploadFile

from .. import db
from ..kb import PRODUCT_CATEGORIES, chunk_text, kb_rows_for_viewer, search_kb_chunks
from ..schemas import KnowledgeUpdate, ProductIn
from ..security import get_current_user, scope_owner
from ..storage import upload_to_storage

log = logging.getLogger("sai")
router = APIRouter(prefix="/api/v1", tags=["catalog"])


# ---------- PRODUCTS ----------
@router.post("/products")
async def create_product(body: ProductIn, cu: dict = Depends(get_current_user)):
    pid = body.id or str(uuid.uuid4())
    data = {"id": pid, "owner_id": cu["id"], "name": body.name,
            "category": body.category if body.category in PRODUCT_CATEGORIES else "Umum",
            "description": body.description, "price_range": body.price_range,
            "duration": body.duration, "kb_text": body.kb_text,
            "created_at": datetime.datetime.now(datetime.timezone.utc).isoformat()}
    db.db_exec("INSERT OR REPLACE INTO products (id,owner_id,name,category,description,price_range,duration,kb_text,created_at) "
               "VALUES (:id,:owner_id,:name,:category,:description,:price_range,:duration,:kb_text,:created_at)", data)
    return {"status": "saved", "product_id": pid}


@router.get("/products")
async def list_products(req: Request, category: str = None, cu: dict = Depends(get_current_user)):
    owner_id = scope_owner(cu, req.query_params.get("owner"))
    if owner_id is None:
        rows = db.db_query("SELECT * FROM products ORDER BY created_at DESC")
    else:
        rows = db.db_query("SELECT * FROM products WHERE owner_id=? ORDER BY created_at DESC", (owner_id,))
    if category:
        rows = [r for r in rows if r.get("category") == category]
    return [dict(r) for r in rows]


@router.get("/products/{pid}")
async def get_product(pid: str, cu: dict = Depends(get_current_user)):
    rows = db.db_query("SELECT * FROM products WHERE id=?", (pid,))
    if not rows:
        raise HTTPException(404, "Produk tidak ditemukan")
    return dict(rows[0])


# ---------- UPLOAD / KNOWLEDGE ----------
@router.post("/upload")
async def upload_product_file(file: UploadFile = File(...), cu: dict = Depends(get_current_user)):
    result = await upload_to_storage(file)
    return result


def _extract_text(content_bytes: bytes, ext: str, filename: str) -> str:
    kb_text = ""
    ext_lower = ext
    try:
        if ext_lower in (".txt", ".md"):
            kb_text = content_bytes.decode("utf-8", errors="replace")
        elif ext_lower == ".pdf":
            try:
                from pypdf import PdfReader
                reader = PdfReader(io.BytesIO(content_bytes))
                kb_text = "\n\n".join((page.extract_text() or "") for page in reader.pages)
            except ImportError:
                import fitz
                doc = fitz.open(stream=content_bytes, filetype="pdf")
                kb_text = "\n\n".join(page.get_text() for page in doc)
                doc.close()
        elif ext_lower == ".docx":
            from docx import Document
            doc = Document(io.BytesIO(content_bytes))
            kb_text = "\n\n".join(p.text or "" for p in doc.paragraphs if p.text.strip())
        elif ext_lower in (".xlsx", ".xls"):
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(content_bytes), read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        rows.append("\t".join(cells))
            kb_text = "\n".join(rows)
            wb.close()
        elif ext_lower == ".pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content_bytes))
            kb_text = "\n\n".join(shape.text or "" for slide in prs.slides
                                  for shape in slide.shapes if hasattr(shape, "text"))
    except ImportError as ie:
        log.warning("[KB Upload] Missing dep for %s: %s", ext, ie)
        kb_text = "(extraction error: dependency tidak tersedia)"
    except Exception:
        log.exception("[KB Upload] %s extraction error", ext)
        kb_text = "(extraction error)"
    if not kb_text:
        kb_text = f"File KB: {filename}. Buka dokumen untuk detail."
    return kb_text


@router.post("/knowledge/upload")
async def upload_knowledge(file: UploadFile = File(...), category: str = Form("Umum"),
                           name: str = Form(""), extract_text: bool = Form(True),
                           cu: dict = Depends(get_current_user)):
    if category not in PRODUCT_CATEGORIES:
        category = "Umum"
    content_bytes = await file.read()
    ext = os.path.splitext(file.filename or "")[1].lower()
    kb_text = _extract_text(content_bytes, ext, file.filename or "") if extract_text \
        else f"File KB: {file.filename}"
    try:
        reuse_file = UploadFile(filename=file.filename, file=io.BytesIO(content_bytes), headers={})
        upload_result = await upload_to_storage(reuse_file)
    except Exception:
        log.exception("[KB Upload] Storage error")
        upload_result = {"status": "error", "error": "storage_error"}
    chunks = chunk_text(kb_text)
    doc_id = str(uuid.uuid4())
    now = datetime.datetime.now(datetime.timezone.utc).isoformat()
    db.db_exec("INSERT INTO knowledge_base (id,owner_id,category,name,filename,file_url,kb_text,chunks,chunk_count,uploaded_at) "
               "VALUES (?,?,?,?,?,?,?,?,?,?)",
               (doc_id, cu["id"], category, name or file.filename or "", file.filename or "",
                upload_result.get("public_url", ""), kb_text[:50000], json.dumps(chunks), len(chunks), now))
    return {"status": "saved", "doc_id": doc_id, "kb_text_length": len(kb_text),
            "chunk_count": len(chunks), "chunks_preview": chunks[:3], "file": upload_result}


@router.get("/knowledge/search")
async def knowledge_search(req: Request, q: str = "", cat: str = None, limit: int = 5,
                           cu: dict = Depends(get_current_user)):
    if not q.strip():
        return {"status": "error", "error": "q wajib"}
    category = None
    if cat and cat.strip() and cat.strip() in PRODUCT_CATEGORIES:
        category = cat.strip()
    owner_filter = scope_owner(cu, req.query_params.get("owner"))
    hits = search_kb_chunks(q, category=category, limit=max(1, min(limit, 20)),
                            cu=cu, owner_filter=owner_filter)
    return {"status": "success", "data": {"query": q, "category": category,
                                          "count": len(hits), "hits": hits}}


def _kb_doc_dict(p: dict) -> dict:
    ch = p.get("chunks") or []
    if isinstance(ch, str):
        try:
            ch = json.loads(ch)
        except Exception:
            ch = []
    return {"id": p.get("id"), "owner_id": p.get("owner_id"), "name": p.get("name"),
            "category": p.get("category"), "filename": p.get("filename", ""),
            "url": p.get("file_url"), "uploaded_at": p.get("uploaded_at"),
            "chunk_count": p.get("chunk_count") or len(ch),
            "preview": (p.get("kb_text") or "")[:400], "kb_text": p.get("kb_text") or ""}


@router.get("/knowledge")
async def list_knowledge(req: Request, category: str = None, cu: dict = Depends(get_current_user)):
    if category and category not in PRODUCT_CATEGORIES:
        return {"status": "error", "error": f"Invalid category. Valid: {PRODUCT_CATEGORIES}"}
    rows = kb_rows_for_viewer(cu, scope_owner(cu, req.query_params.get("owner")))
    if category:
        rows = [r for r in rows if r.get("category") == category]
    docs = [_kb_doc_dict(r) for r in rows]
    counts = {c: 0 for c in PRODUCT_CATEGORIES}
    for d in docs:
        counts[d.get("category")] = counts.get(d.get("category"), 0) + 1
    total_chunks = sum(d.get("chunk_count") or 0 for d in docs)
    return {"status": "success", "data": {"category": category, "count": len(docs), "docs": docs,
                                          "category_counts": counts, "total_chunks": total_chunks}}


def _kb_doc_or_error(doc_id: str, cu: dict):
    rows = db.db_query("SELECT * FROM knowledge_base WHERE id=?", (doc_id,))
    if not rows:
        raise HTTPException(404, "Dokumen tidak ditemukan")
    r = rows[0]
    if cu["role"] != "admin" and r["owner_id"] not in (cu["id"], "*"):
        raise HTTPException(403, "Tidak punya akses")
    return r


@router.get("/knowledge/doc/{doc_id}")
async def get_knowledge_doc(doc_id: str, cu: dict = Depends(get_current_user)):
    r = _kb_doc_or_error(doc_id, cu)
    return {"status": "success", "data": _kb_doc_dict(r)}


@router.put("/knowledge/{doc_id}")
async def update_knowledge(doc_id: str, body: KnowledgeUpdate, cu: dict = Depends(get_current_user)):
    _kb_doc_or_error(doc_id, cu)
    name = (body.name or "").strip()
    category = (body.category or "").strip()
    kb_text = (body.kb_text or "").strip()
    if category and category not in PRODUCT_CATEGORIES:
        return {"status": "error", "error": f"Invalid category. Valid: {PRODUCT_CATEGORIES}"}
    if not name and not kb_text:
        return {"status": "error", "error": "name atau kb_text wajib"}
    chunks = chunk_text(kb_text) if kb_text else None
    sets, params = [], []
    if name:
        sets.append("name=?")
        params.append(name)
    if category:
        sets.append("category=?")
        params.append(category)
    if kb_text:
        sets.append("kb_text=?")
        params.append(kb_text[:50000])
        sets.append("chunks=?")
        params.append(json.dumps(chunks))
        sets.append("chunk_count=?")
        params.append(len(chunks))
    params.append(doc_id)
    db.db_exec(f"UPDATE knowledge_base SET {','.join(sets)} WHERE id=?", tuple(params))
    return {"status": "saved", "id": doc_id}


@router.delete("/knowledge/{doc_id}")
async def delete_knowledge(doc_id: str, cu: dict = Depends(get_current_user)):
    _kb_doc_or_error(doc_id, cu)
    db.db_exec("DELETE FROM knowledge_base WHERE id=?", (doc_id,))
    return {"status": "deleted", "id": doc_id}


@router.get("/knowledge/{category}")
async def get_knowledge(category: str, req: Request, cu: dict = Depends(get_current_user)):
    if category not in PRODUCT_CATEGORIES:
        return {"status": "error", "error": f"Invalid category. Valid: {PRODUCT_CATEGORIES}"}
    rows = kb_rows_for_viewer(cu, scope_owner(cu, req.query_params.get("owner")))
    rows = [r for r in rows if r.get("category") == category]
    docs = []
    for p in rows:
        ch = p.get("chunks") or []
        if isinstance(ch, str):
            try:
                ch = json.loads(ch)
            except Exception:
                ch = []
        docs.append({"id": p.get("id"), "name": p.get("name"), "category": p.get("category"),
                     "url": p.get("file_url"), "uploaded_at": p.get("uploaded_at"),
                     "chunk_count": p.get("chunk_count") or len(ch),
                     "preview": (p.get("kb_text") or "")[:400]})
    return {"category": category, "docs": docs, "count": len(docs)}
